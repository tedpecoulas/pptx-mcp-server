"""
Claude Skills MCP Server - PPTX Edition with SSE Support
Python server for reading and modifying PowerPoint templates
WITH INTELLIGENT FONT AUTO-SIZING v2.3 - Dual Groups + No double bullets
"""

from flask import Flask, request, jsonify, send_file, Response
from flask_cors import CORS
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt, Inches
from pptx.enum.text import MSO_AUTO_SIZE, PP_PARAGRAPH_ALIGNMENT
import requests
import io
import json
import tempfile
import os
import time
from datetime import datetime
import re
import math

app = Flask(__name__)
CORS(app)

# Store modified presentations temporarily
temp_files = {}

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# CONFIGURATION DES GROUPES DE FORMATTAGE
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# GROUPE 1 : Même taille de police pour ces 3 cadres
GROUP_1_SHAPES = ["contexte", "résultats", "travaux réalisés"]

# GROUPE 2 : Même taille de police pour ces 2 cadres (indépendante du groupe 1)
GROUP_2_SHAPES = ["type de mission", "outils utilisés"]

# Shapes qui ne doivent PAS avoir de bullets (texte en paragraphes)
NO_BULLET_SHAPES = ["contexte"]

# Taille de police
DEFAULT_FONT_SIZE = 12
MIN_FONT_SIZE = 8  # Taille minimale absolue
MAX_FONT_SIZE = 14

# Interligne pour esthétique
LINE_SPACING = 1.2


def sanitize_filename(text):
    """Sanitize text for use in filename"""
    text = re.sub(r'[<>:"/\\|?*]', '-', text)
    text = text.strip(' .')
    return text[:50] if text else "Document"


def download_pptx(url):
    """Download PPTX from URL and return Presentation object"""
    response = requests.get(url, timeout=30)
    response.raise_for_status()
    pptx_bytes = io.BytesIO(response.content)
    return Presentation(pptx_bytes)


def normalize_shape_name(name):
    """Normalise le nom d'une shape pour comparaison"""
    return name.lower().strip()


def get_shape_group(shape):
    """
    Détermine à quel groupe appartient une shape
    Retourne 1 (Contexte/Résultats/Travaux), 2 (Type/Outils), ou None
    """
    if not shape.has_text_frame:
        return None
    
    shape_name_normalized = normalize_shape_name(shape.name)
    shape_text_normalized = normalize_shape_name(shape.text_frame.text) if shape.text_frame.text else ""
    
    # Vérifier Groupe 1 : Contexte, Résultats, Travaux réalisés
    for keyword in GROUP_1_SHAPES:
        if keyword.lower() in shape_name_normalized or keyword.lower() in shape_text_normalized:
            return 1
    
    # Vérifier Groupe 2 : Type de mission, Outils utilisés
    for keyword in GROUP_2_SHAPES:
        if keyword.lower() in shape_name_normalized or keyword.lower() in shape_text_normalized:
            return 2
    
    return None


def should_have_bullets(shape):
    """
    Détermine si une shape doit avoir des bullet points
    Retourne False pour "Contexte", True pour les autres
    """
    if not shape.has_text_frame:
        return False
    
    shape_name_normalized = normalize_shape_name(shape.name)
    shape_text_normalized = normalize_shape_name(shape.text_frame.text) if shape.text_frame.text else ""
    
    # Vérifier si c'est une shape "no bullets"
    for keyword in NO_BULLET_SHAPES:
        if keyword.lower() in shape_name_normalized or keyword.lower() in shape_text_normalized:
            return False
    
    return True


def estimate_text_height(text, font_size, shape_width, line_spacing=1.2):
    """
    Estime la hauteur du texte rendu en fonction de la longueur, 
    taille de police, largeur du cadre et interligne
    """
    # Estimation du nombre de caractères par ligne
    chars_per_inch = 72 / (font_size * 0.5)
    shape_width_points = shape_width.inches * 72
    
    # Réduire la largeur effective pour tenir compte des marges
    effective_width = shape_width_points * 0.9  # 10% de marge
    chars_per_line = effective_width / (font_size * 0.5)
    
    # Calculer le nombre de lignes
    text_length = len(text)
    explicit_lines = text.count('\n') + 1
    wrapped_lines = math.ceil(text_length / chars_per_line)
    total_lines = max(explicit_lines, wrapped_lines)
    
    # Hauteur totale
    line_height_points = font_size * line_spacing
    total_height_points = total_lines * line_height_points
    total_height_inches = total_height_points / 72
    
    return total_height_inches, total_lines


def find_optimal_font_size(texts_and_shapes, max_size=DEFAULT_FONT_SIZE, min_size=MIN_FONT_SIZE, line_spacing=1.2):
    """
    Trouve la taille de police optimale pour un GROUPE de shapes
    La taille sera la MÊME pour toutes les shapes du groupe
    Calcule selon la shape la plus contraignante (texte le plus long / cadre le plus petit)
    """
    if not texts_and_shapes:
        return max_size
    
    optimal_size = max_size
    
    # Tester chaque shape du groupe pour trouver la taille max qui fonctionne pour TOUTES
    for text, shape in texts_and_shapes:
        if not text or not shape.has_text_frame:
            continue
        
        shape_height = shape.height
        shape_width = shape.width
        
        # Tester différentes tailles de max_size à min_size
        for test_size in range(max_size, min_size - 1, -1):
            estimated_height, num_lines = estimate_text_height(
                text, test_size, shape_width, line_spacing
            )
            
            # Marge de sécurité de 15%
            safety_margin = shape_height.inches * 0.15
            available_height = shape_height.inches - safety_margin
            
            if estimated_height <= available_height:
                # Cette taille convient pour CETTE shape
                # On prend le minimum pour que ça fonctionne pour TOUTES les shapes du groupe
                optimal_size = min(optimal_size, test_size)
                print(f"  📐 '{shape.name}': {len(text)}c, {num_lines}L → {test_size}pt OK (hauteur: {shape_height.inches:.2f}\")")
                break
        else:
            # Aucune taille ne convient, utiliser la taille minimale
            optimal_size = min_size
            print(f"  ⚠️ '{shape.name}': Texte trop long ({len(text)}c), taille min {min_size}pt forcée")
    
    return optimal_size


def clean_bullet_text(text):
    """
    Nettoie le texte en enlevant les bullets du texte lui-même
    car PowerPoint les ajoutera automatiquement via son système de bullets
    """
    if not text:
        return text
    
    lines = text.split('\n')
    cleaned_lines = []
    
    for line in lines:
        line = line.strip()
        if not line:
            continue
        
        # Enlever tous les types de bullets du texte
        if line.startswith('• '):
            line = line[2:]
        elif line.startswith('•'):
            line = line[1:].strip()
        elif line.startswith('- '):
            line = line[2:]
        elif line.startswith('-'):
            line = line[1:].strip()
        elif line.startswith('* '):
            line = line[2:]
        elif line.startswith('*'):
            line = line[1:].strip()
        
        cleaned_lines.append(line)
    
    return '\n'.join(cleaned_lines)


def apply_text_with_formatting(shape, text, font_size, line_spacing=1.2, use_bullets=True):
    """
    Applique le texte avec formatage optimisé
    - font_size : Taille de police (définie par groupe)
    - line_spacing : Interligne
    - use_bullets : False pour "Contexte" (paragraphes), True pour autres (bullets)
    """
    if not shape.has_text_frame:
        return False
    
    # Déterminer si on doit utiliser les bullets
    should_use_bullets = use_bullets and should_have_bullets(shape)
    
    # Nettoyer le texte des bullets existants si on va les réactiver
    cleaned_text = clean_bullet_text(text) if should_use_bullets else text
    
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    
    # Réduire les marges internes pour maximiser l'espace disponible
    text_frame.margin_bottom = Inches(0.05)
    text_frame.margin_top = Inches(0.05)
    text_frame.margin_left = Inches(0.1)
    text_frame.margin_right = Inches(0.1)
    
    # Séparer les lignes
    lines = cleaned_text.split('\n')
    
    for i, line in enumerate(lines):
        if not line.strip():
            continue
            
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = line
        p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
        p.line_spacing = line_spacing
        p.level = 0
        
        # Espacement selon le type
        if should_use_bullets:
            # Bullets : PowerPoint les ajoutera automatiquement
            p.space_before = Pt(2)
            p.space_after = Pt(2)
        else:
            # Paragraphes (Contexte) : pas de bullets
            p.space_before = Pt(0)
            p.space_after = Pt(4)
        
        # Appliquer la taille de police
        for run in p.runs:
            run.font.size = Pt(font_size)
    
    bullet_status = "bullets" if should_use_bullets else "paragraphes"
    print(f"  ✍️  '{shape.name}': {len(text)}c, {len(lines)}L → {font_size}pt ({bullet_status})")
    return True


def analyze_presentation(prs):
    """Analyze presentation structure and return JSON"""
    analysis = {
        "total_slides": len(prs.slides),
        "slides": []
    }
    
    for slide_idx, slide in enumerate(prs.slides):
        slide_info = {
            "slide_number": slide_idx,
            "layout_name": slide.slide_layout.name,
            "shapes": []
        }
        
        for shape_idx, shape in enumerate(slide.shapes):
            shape_info = {
                "shape_id": shape_idx,
                "name": shape.name,
                "type": str(shape.shape_type),
                "has_text_frame": shape.has_text_frame,
                "group": get_shape_group(shape),
                "should_have_bullets": should_have_bullets(shape)
            }
            
            if shape.has_text_frame:
                text = shape.text_frame.text
                shape_info["text"] = text
                shape_info["text_length"] = len(text)
                shape_info["width_inches"] = round(shape.width.inches, 2)
                shape_info["height_inches"] = round(shape.height.inches, 2)
                
                if shape.is_placeholder:
                    shape_info["placeholder_type"] = str(shape.placeholder_format.type)
                else:
                    shape_info["placeholder_type"] = None
                
                shape_info["paragraph_count"] = len(shape.text_frame.paragraphs)
            
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                shape_info["is_picture"] = True
            
            slide_info["shapes"].append(shape_info)
        
        analysis["slides"].append(slide_info)
    
    return analysis


def modify_presentation(prs, modifications):
    """
    Modifie la présentation avec ajustement intelligent de la police en 2 GROUPES
    
    GROUPE 1 (Contexte, Résultats, Travaux) : Même taille de police pour les 3
    GROUPE 2 (Type de mission, Outils) : Même taille de police pour les 2 (indépendante du groupe 1)
    """
    warnings = []
    
    print("\n" + "="*70)
    print("  PHASE 1 : COLLECTE DES SHAPES PAR GROUPE")
    print("="*70)
    
    # Phase 1 : Collecter les shapes par groupe
    group_1_data = []  # Contexte, Résultats, Travaux réalisés
    group_2_data = []  # Type de mission, Outils utilisés
    other_shapes_data = []
    
    for slide_key, shape_mods in modifications.items():
        slide_num = int(slide_key.split('_')[1])
        
        if slide_num >= len(prs.slides):
            continue
        
        slide = prs.slides[slide_num]
        
        for shape_key, new_text in shape_mods.items():
            shape_num = int(shape_key.split('_')[1])
            
            if shape_num >= len(slide.shapes):
                continue
            
            shape = slide.shapes[shape_num]
            group = get_shape_group(shape)
            
            if group == 1:
                group_1_data.append((new_text, shape))
                print(f"  ✓ Groupe 1: {shape.name} ({len(new_text)} caractères)")
            elif group == 2:
                group_2_data.append((new_text, shape))
                print(f"  ✓ Groupe 2: {shape.name} ({len(new_text)} caractères)")
            else:
                other_shapes_data.append((new_text, shape))
                print(f"  ✓ Autre: {shape.name} ({len(new_text)} caractères)")
    
    print("\n" + "="*70)
    print("  PHASE 2 : CALCUL DES TAILLES OPTIMALES PAR GROUPE")
    print("="*70)
    
    # Phase 2 : Calculer la taille optimale pour GROUPE 1
    print(f"\n🎯 [GROUPE 1] {len(group_1_data)} shapes (Contexte, Résultats, Travaux)")
    print("   → Taille de police UNIFORME pour les 3 cadres")
    group_1_font_size = DEFAULT_FONT_SIZE
    
    if group_1_data:
        group_1_font_size = find_optimal_font_size(
            group_1_data, 
            max_size=MAX_FONT_SIZE, 
            min_size=MIN_FONT_SIZE, 
            line_spacing=LINE_SPACING
        )
        print(f"\n  ✅ Taille finale Groupe 1 : {group_1_font_size}pt (pour les 3 cadres)\n")
        
        if group_1_font_size == MIN_FONT_SIZE:
            warnings.append(
                f"⚠️ GROUPE 1 (Contexte, Résultats, Travaux) : Le texte est très dense. "
                f"La police a été réduite au minimum ({MIN_FONT_SIZE}pt). "
                f"Pour améliorer la lisibilité, réduisez le contenu de ces cadres."
            )
    
    # Phase 2bis : Calculer la taille optimale pour GROUPE 2 (INDÉPENDANTE du groupe 1)
    print(f"🎯 [GROUPE 2] {len(group_2_data)} shapes (Type de mission, Outils)")
    print("   → Taille de police UNIFORME pour les 2 cadres (indépendante du Groupe 1)")
    group_2_font_size = DEFAULT_FONT_SIZE
    
    if group_2_data:
        group_2_font_size = find_optimal_font_size(
            group_2_data, 
            max_size=MAX_FONT_SIZE, 
            min_size=MIN_FONT_SIZE, 
            line_spacing=LINE_SPACING
        )
        print(f"\n  ✅ Taille finale Groupe 2 : {group_2_font_size}pt (pour les 2 cadres)\n")
        
        if group_2_font_size == MIN_FONT_SIZE:
            warnings.append(
                f"⚠️ GROUPE 2 (Type de mission, Outils) : Le texte est très dense. "
                f"La police a été réduite au minimum ({MIN_FONT_SIZE}pt). "
                f"Pour améliorer la lisibilité, réduisez le contenu de ces cadres."
            )
    
    print("="*70)
    print("  PHASE 3 : APPLICATION DES MODIFICATIONS")
    print("="*70 + "\n")
    
    # Phase 3 : Appliquer les modifications avec les tailles calculées
    print(f"📝 Application Groupe 1 ({group_1_font_size}pt) :")
    for text, shape in group_1_data:
        use_bullets = should_have_bullets(shape)
        apply_text_with_formatting(shape, text, group_1_font_size, LINE_SPACING, use_bullets=use_bullets)
    
    print(f"\n📝 Application Groupe 2 ({group_2_font_size}pt) :")
    for text, shape in group_2_data:
        apply_text_with_formatting(shape, text, group_2_font_size, LINE_SPACING, use_bullets=True)
    
    print(f"\n📝 Application Autres shapes :")
    for text, shape in other_shapes_data:
        individual_size = find_optimal_font_size([(text, shape)], max_size=MAX_FONT_SIZE, min_size=MIN_FONT_SIZE, line_spacing=1.0)
        apply_text_with_formatting(shape, text, individual_size, 1.0, use_bullets=True)
    
    print("\n" + "="*70)
    print("  ✅ GÉNÉRATION TERMINÉE")
    print("="*70 + "\n")
    
    return prs, warnings


def handle_mcp_request(body, request_id):
    """Handle MCP JSON-RPC request and return response"""
    method = body.get('method', '')
    params = body.get('params', {})
    
    print(f"📥 Method: {method}")
    
    # Route: initialize
    if method == 'initialize':
        client_version = params.get('protocolVersion', '2025-06-18')
        return {
            "jsonrpc": "2.0",
            "id": request_id,
            "result": {
                "protocolVersion": client_version,
                "capabilities": {
                    "tools": {"listChanged": False},
                    "resources": {},
                    "prompts": {}
                },
                "serverInfo": {
                    "name": "pptx-mcp-server",
                    "version": "2.3.0"
                }
            }
        }
    
    # Route: tools/list
    if method == 'tools/list':
        return {
            "jsonrpc": "2.0",
            "id": request_id,
            "result": {
                "tools": [
                    {
                        "name": "analyze_template",
                        "description": "Analyse la structure d'un template PowerPoint",
                        "inputSchema": {
                            "type": "object",
                            "properties": {
                                "template_url": {
                                    "type": "string",
                                    "description": "URL du fichier PPTX à analyser"
                                }
                            },
                            "required": ["template_url"]
                        }
                    },
                    {
                        "name": "modify_template",
                        "description": "Modifie un template PowerPoint avec sizing uniforme par groupe (Groupe 1: Contexte/Résultats/Travaux, Groupe 2: Type/Outils)",
                        "inputSchema": {
                            "type": "object",
                            "properties": {
                                "template_url": {
                                    "type": "string",
                                    "description": "URL du template PPTX"
                                },
                                "modifications": {
                                    "type": "object",
                                    "description": "Dictionnaire des modifications"
                                },
                                "metadata": {
                                    "type": "object",
                                    "description": "Métadonnées pour nommer le fichier",
                                    "properties": {
                                        "client": {"type": "string"},
                                        "mission": {"type": "string"},
                                        "consultant": {"type": "string"}
                                    }
                                }
                            },
                            "required": ["template_url", "modifications"]
                        }
                    }
                ]
            }
        }
    
    # Route: tools/call
    if method == 'tools/call':
        tool_name = params.get('name')
        args = params.get('arguments', {})
        
        if tool_name == 'analyze_template':
            try:
                template_url = args.get('template_url')
                print(f"📄 Analyzing template: {template_url}")
                
                prs = download_pptx(template_url)
                analysis = analyze_presentation(prs)
                
                return {
                    "jsonrpc": "2.0",
                    "id": request_id,
                    "result": {
                        "content": [{
                            "type": "text",
                            "text": json.dumps(analysis, indent=2, ensure_ascii=False)
                        }]
                    }
                }
            except Exception as e:
                return {
                    "jsonrpc": "2.0",
                    "id": request_id,
                    "error": {
                        "code": -32603,
                        "message": f"Error analyzing template: {str(e)}"
                    }
                }
        
        if tool_name == 'modify_template':
            try:
                template_url = args.get('template_url')
                modifications = args.get('modifications', {})
                metadata = args.get('metadata', {})
                
                print(f"✏️ Modifying template: {template_url}")
                print(f"✏️ Metadata: {metadata}")
                
                prs = download_pptx(template_url)
                prs, warnings = modify_presentation(prs, modifications)
                
                client = sanitize_filename(metadata.get('client', ''))
                mission = sanitize_filename(metadata.get('mission', ''))
                consultant = sanitize_filename(metadata.get('consultant', ''))
                
                timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
                file_id = f"pptx_{timestamp}"
                
                if client and mission and consultant:
                    suggested_name = f"REX - {client} - {mission} - {consultant}.pptx"
                elif client and mission:
                    suggested_name = f"REX - {client} - {mission}.pptx"
                elif client:
                    suggested_name = f"REX - {client}.pptx"
                else:
                    suggested_name = f"REX_{timestamp}.pptx"
                
                output_file = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
                prs.save(output_file.name)
                
                temp_files[file_id] = {
                    'path': output_file.name,
                    'suggested_name': suggested_name
                }
                
                base_url = os.environ.get('SERVER_URL', 'https://pptx-mcp-server-production.up.railway.app')
                download_url = f"{base_url}/download/{file_id}"
                
                response_text = f"✅ Votre REX est prêt !\n\n📥 Télécharger ici: {download_url}\n\n💡 Nom de fichier: {suggested_name}\n\n"
                
                if warnings:
                    response_text += "\n" + "\n\n".join(warnings)
                
                return {
                    "jsonrpc": "2.0",
                    "id": request_id,
                    "result": {
                        "content": [{
                            "type": "text",
                            "text": response_text
                        }]
                    }
                }
            except Exception as e:
                print(f"❌ Error: {str(e)}")
                import traceback
                traceback.print_exc()
                return {
                    "jsonrpc": "2.0",
                    "id": request_id,
                    "error": {
                        "code": -32603,
                        "message": f"Error modifying template: {str(e)}"
                    }
                }
        
        return {
            "jsonrpc": "2.0",
            "id": request_id,
            "error": {
                "code": -32601,
                "message": f"Unknown tool: {tool_name}"
            }
        }
    
    return {
        "jsonrpc": "2.0",
        "id": request_id,
        "error": {
            "code": -32601,
            "message": f"Method not found: {method}"
        }
    }


@app.route('/api/mcp', methods=['GET', 'POST', 'OPTIONS'])
def mcp_endpoint():
    """Main MCP endpoint"""
    
    if request.method == 'OPTIONS':
        return '', 200
    
    if request.method == 'GET':
        return jsonify({
            "name": "PPTX MCP Server",
            "version": "2.3.0",
            "tools": ["analyze_template", "modify_template"],
            "groups": {
                "group_1": GROUP_1_SHAPES,
                "group_2": GROUP_2_SHAPES
            }
        })
    
    accept_header = request.headers.get('Accept', '')
    wants_sse = 'text/event-stream' in accept_header
    
    body = request.get_json() or {}
    request_id = body.get('id', 1)
    
    if wants_sse:
        def generate_sse():
            response_data = handle_mcp_request(body, request_id)
            sse_data = f"data: {json.dumps(response_data)}\n\n"
            yield sse_data
            time.sleep(0.5)
        
        return Response(
            generate_sse(),
            mimetype='text/event-stream',
            headers={
                'Cache-Control': 'no-cache',
                'X-Accel-Buffering': 'no',
                'Connection': 'keep-alive'
            }
        )
    
    response_data = handle_mcp_request(body, request_id)
    return jsonify(response_data)


@app.route('/download/<file_id>')
def download_file(file_id):
    """Download endpoint"""
    if file_id not in temp_files:
        return jsonify({"error": "File not found"}), 404
    
    file_info = temp_files[file_id]
    file_path = file_info['path']
    suggested_name = file_info['suggested_name']
    
    if not os.path.exists(file_path):
        return jsonify({"error": "File no longer exists"}), 404
    
    return send_file(
        file_path,
        mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
        as_attachment=True,
        download_name=suggested_name
    )


@app.route('/health')
def health():
    """Health check"""
    return jsonify({
        "status": "healthy",
        "server": "pptx-mcp-server",
        "version": "2.3.0",
        "features": {
            "dual_group_sizing": True,
            "group_1": GROUP_1_SHAPES,
            "group_2": GROUP_2_SHAPES,
            "min_font_size": MIN_FONT_SIZE
        }
    })


if __name__ == '__main__':
    port = int(os.environ.get('PORT', 5000))
    app.run(host='0.0.0.0', port=port, debug=True)